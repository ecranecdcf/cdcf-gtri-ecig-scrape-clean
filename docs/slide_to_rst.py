from pptx import Presentation
import os

def pptx_to_rst(pptx_path, rst_path):
    """
    Converts a PowerPoint presentation (PPTX) to a reStructuredText (RST) file.
    This is a basic example and may require significant customization
    for complex presentations.
    """
    prs = Presentation(pptx_path)
    rst_content = []

    for i, slide in enumerate(prs.slides):
        rst_content.append(f"Slide {i + 1}")
        rst_content.append("=" * len(f"Slide {i + 1}")) # Underline for section title
        rst_content.append("\n")

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Basic text extraction
                rst_content.append(shape.text)
                rst_content.append("\n")
            # You might need to handle other shape types (images, tables, etc.)
            # and convert them to appropriate RST directives.

        rst_content.append("\n") # Add a blank line between slides

    with open(rst_path, "w", encoding="utf-8") as f:
        f.write("\n".join(rst_content))

# Example usage:
# pptx_to_rst("my_presentation.pptx", "output.rst")

if __name__ == "__main__":
    pptx_files = [f for f in os.listdir("../Downloads/slides") if f.endswith(".pptx")]
    print(pptx_files)
    for m in pptx_files:
        pptx_path = os.path.join("../Downloads/slides", m)
        f_name = f"slides_{os.path.splitext(m)[0]}.rst"
        rst_path = os.path.join("docs", f_name)
        pptx_to_rst(pptx_path, rst_path)

        print(f_name)